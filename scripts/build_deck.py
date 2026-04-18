"""Build executive presentation deck (20 min, simple design).

Run:
    python scripts/build_deck.py

Output:
    docs/presentation.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ------------------------------------------------------------------
# Palette (simple, executive-friendly)
# ------------------------------------------------------------------
INK = RGBColor(0x18, 0x18, 0x1B)       # near-black (zinc-900)
BODY = RGBColor(0x3F, 0x3F, 0x46)      # zinc-700
MUTE = RGBColor(0x71, 0x71, 0x7A)      # zinc-500
LINE = RGBColor(0xE4, 0xE4, 0xE7)      # zinc-200
BG   = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x1F, 0x2A, 0x44)    # deep navy (단색 포인트)
ACCENT_SOFT = RGBColor(0xEE, 0xF2, 0xF7)

FONT = "Calibri"

# 16:9 widescreen
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


# ------------------------------------------------------------------
# Low-level helpers
# ------------------------------------------------------------------

def set_bg(slide, rgb=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_text(
    slide,
    left, top, width, height,
    text,
    *,
    size=14,
    bold=False,
    color=BODY,
    align=PP_ALIGN.LEFT,
    font=FONT,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(
    slide,
    left, top, width, height,
    items,
    *,
    size=16,
    color=BODY,
    bullet="•",
    line_spacing=1.35,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_hline(slide, left, top, width, color=LINE, weight=0.75):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_rect(slide, left, top, width, height, fill=ACCENT_SOFT, line=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(0.5)
    rect.shadow.inherit = False
    return rect


def add_footer(slide, page_num, total):
    add_text(
        slide,
        Inches(0.5), Inches(SLIDE_H_IN - 0.35), Inches(8), Inches(0.25),
        "Game Competitor Analysis PoC  ·  2026-04-18",
        size=9, color=MUTE,
    )
    add_text(
        slide,
        Inches(SLIDE_W_IN - 1.5), Inches(SLIDE_H_IN - 0.35), Inches(1.0), Inches(0.25),
        f"{page_num} / {total}",
        size=9, color=MUTE, align=PP_ALIGN.RIGHT,
    )


def add_title(slide, text, eyebrow=None):
    top = Inches(0.55)
    if eyebrow:
        add_text(
            slide,
            Inches(0.6), top, Inches(12), Inches(0.3),
            eyebrow.upper(),
            size=10, bold=True, color=ACCENT,
        )
        top = Inches(0.9)
    add_text(
        slide,
        Inches(0.6), top, Inches(12.1), Inches(0.7),
        text,
        size=30, bold=True, color=INK,
    )
    add_hline(slide, Inches(0.6), Inches(1.75), Inches(12.1))


# ------------------------------------------------------------------
# Slide builders
# ------------------------------------------------------------------

def build(prs, page_total_provider):
    # python-pptx needs us to know total ahead of footer; we count slides after build.
    slides = []

    def new_slide():
        layout = prs.slide_layouts[6]  # blank
        s = prs.slides.add_slide(layout)
        set_bg(s)
        slides.append(s)
        return s

    # ---- 1. Cover ----------------------------------------------------
    s = new_slide()
    add_rect(s, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(2.8), fill=ACCENT)
    add_text(
        s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
        "INTERNAL PM TOOL · POC REVIEW",
        size=11, bold=True, color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    add_text(
        s, Inches(0.8), Inches(1.3), Inches(12), Inches(1.2),
        "Game Competitor Analysis",
        size=44, bold=True, color=BG,
    )
    add_text(
        s, Inches(0.8), Inches(2.1), Inches(12), Inches(0.5),
        "주간 경쟁작 리서치 자동화 시스템",
        size=20, color=RGBColor(0xE4, 0xE8, 0xF0),
    )
    add_text(
        s, Inches(0.8), Inches(3.4), Inches(12), Inches(0.4),
        "경영진 대상 발표 · 20분",
        size=14, bold=True, color=INK,
    )
    add_text(
        s, Inches(0.8), Inches(3.9), Inches(12), Inches(0.4),
        "2026년 4월 18일",
        size=13, color=BODY,
    )
    add_text(
        s, Inches(0.8), Inches(5.9), Inches(12), Inches(0.4),
        "발표 구조: 문제 → 솔루션 → 데모 → 지표 → 로드맵 → 의사결정",
        size=11, color=MUTE,
    )

    # ---- 2. 핵심 메시지 (Exec summary) -------------------------------
    s = new_slide()
    add_title(s, "한 장 요약", eyebrow="Executive Summary")
    add_rect(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.3), fill=ACCENT_SOFT)
    add_text(
        s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(0.9),
        "PM 이 매주 수시간씩 수작업으로 하던 경쟁작 리서치를,\n"
        "4축 유사도 엔진 + 주간 배치 + 내부 웹툴로 자동화했습니다.",
        size=20, bold=True, color=INK,
    )
    add_bullets(
        s, Inches(0.6), Inches(3.8), Inches(12.1), Inches(3.5),
        [
            "담당 게임을 등록하면 매주 월요일 Top 10 경쟁작 리포트가 자동 생성",
            "LLM + 임베딩으로 장르·BM·세션길이·리뷰 기반 경쟁구도를 설명가능하게 수치화",
            "PM 은 리포트에 👍/👎 만 남기면 됨 → 다음 주 가중치에 자동 반영",
            "4주 PoC 로 End-to-End 동작 · 운영비 사실상 $0 (Groq 무료 tier + 로컬 임베딩)",
            "Steam 전용 · 인증 없음 등 범위 제한은 §9 참조, 후속 로드맵은 §11",
        ],
        size=16,
    )

    # ---- 3. 문제 (Why now) -------------------------------------------
    s = new_slide()
    add_title(s, "지금 왜 이게 필요한가", eyebrow="Problem")
    add_rect(s, Inches(0.6), Inches(2.1), Inches(5.9), Inches(4.8),
             fill=RGBColor(0xFE, 0xF2, 0xF2),
             line=RGBColor(0xFE, 0xCA, 0xCA))
    add_text(s, Inches(0.85), Inches(2.3), Inches(5.5), Inches(0.4),
             "현재 상태", size=12, bold=True, color=RGBColor(0x99, 0x1B, 0x1B))
    add_bullets(
        s, Inches(0.85), Inches(2.75), Inches(5.5), Inches(4.0),
        [
            "PM 이 매주 월요일 아침 수동으로 경쟁작 조사",
            "스토어 검색 → 장르·BM·리뷰 훑기 → 스프레드시트 정리",
            "담당 게임 수 × 수 시간 → 월요일이 통째로 녹음",
            "조사 품질이 개인 경험에 의존 → 재현성·비교성 낮음",
            "일회성 메모로 끝나고 기관 기억으로 축적되지 않음",
        ],
        size=14, color=RGBColor(0x7F, 0x1D, 0x1D),
    )
    add_rect(s, Inches(6.8), Inches(2.1), Inches(5.9), Inches(4.8),
             fill=RGBColor(0xEC, 0xFD, 0xF5),
             line=RGBColor(0xA7, 0xF3, 0xD0))
    add_text(s, Inches(7.05), Inches(2.3), Inches(5.5), Inches(0.4),
             "원하는 상태", size=12, bold=True, color=RGBColor(0x06, 0x5F, 0x46))
    add_bullets(
        s, Inches(7.05), Inches(2.75), Inches(5.5), Inches(4.0),
        [
            "월요일 아침 10분 안에 경쟁 구도 파악",
            "시스템이 후보 Top N 제시 → PM 이 취사선택",
            "피드백이 다음 주 품질 개선으로 되돌아옴",
            "주간 스냅샷이 DB 에 누적 → 시장 추이 추적",
            "담당 게임 늘어도 PM 공수는 선형 증가 X",
        ],
        size=14, color=RGBColor(0x06, 0x5F, 0x46),
    )

    # ---- 4. 성공 지표 정의 -------------------------------------------
    s = new_slide()
    add_title(s, "성공은 이렇게 측정합니다", eyebrow="How We Measure Success")
    add_text(
        s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.5),
        "북극성 지표는 \"정확도\" 가 아니라 \"PM 리포트 준비 시간\" 입니다.",
        size=16, bold=True, color=INK,
    )
    add_bullets(
        s, Inches(0.6), Inches(2.8), Inches(12.1), Inches(3.5),
        [
            "북극성: 주간 리포트 준비 시간  X → 1시간 이내",
            "보조: Top 10 precision (gold set 대비) ≥ 0.6 — 후보 제시 수준이면 충분",
            "보조: PM upvote rate ≥ 50% — 유용성의 implicit signal",
            "품질: Feature self-consistency ≥ 90% — LLM 출력 안정성",
            "신뢰성: End-to-end 주간 실행 성공률 ≥ 95%",
        ],
        size=16,
    )
    add_rect(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.9), fill=ACCENT_SOFT)
    add_text(
        s, Inches(0.85), Inches(6.25), Inches(11.5), Inches(0.6),
        "핵심: 경쟁작을 \"확정\" 할 필요 없음. 리포트의 후보를 제시하면 PM 이 판단.",
        size=14, bold=True, color=ACCENT,
    )

    # ---- 5. 경쟁작 정의 (4축) ----------------------------------------
    s = new_slide()
    add_title(s, "'경쟁작' 을 수치로 정의합니다", eyebrow="Competitor Definition")
    add_text(
        s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.5),
        "본질: target user 층의 겹침. PoC 는 4개 축의 연속값 유사도로 근사.",
        size=14, color=BODY,
    )

    # 4 axes table
    col_x = [0.6, 4.2, 8.6, 11.2]
    headers = ["축", "표현 방식", "가중치(초기)"]
    rows = [
        ("Semantic", "description + review 임베딩 cosine", "0.40"),
        ("Genre",    "LLM 추출 장르 임베딩 cosine (이진 X)", "0.25"),
        ("Tier",     "log1p(review_count) 정규화 Gaussian", "0.20"),
        ("BM",       "{gacha, ads, premium, sub} KL 대칭거리", "0.15"),
    ]
    head_top = Inches(2.7)
    add_rect(s, Inches(0.6), head_top, Inches(12.1), Inches(0.5), fill=ACCENT)
    add_text(s, Inches(0.85), head_top + Emu(40000), Inches(3.0), Inches(0.4),
             headers[0], size=12, bold=True, color=BG)
    add_text(s, Inches(4.2),  head_top + Emu(40000), Inches(7.0), Inches(0.4),
             headers[1], size=12, bold=True, color=BG)
    add_text(s, Inches(11.2), head_top + Emu(40000), Inches(1.5), Inches(0.4),
             headers[2], size=12, bold=True, color=BG)

    row_h = 0.55
    for i, (axis, expr, weight) in enumerate(rows):
        y = 3.2 + i * row_h
        if i % 2 == 1:
            add_rect(s, Inches(0.6), Inches(y), Inches(12.1), Inches(row_h),
                     fill=RGBColor(0xF8, 0xFA, 0xFC))
        add_text(s, Inches(0.85), Inches(y + 0.1), Inches(3.0), Inches(0.4),
                 axis, size=13, bold=True, color=INK)
        add_text(s, Inches(4.2), Inches(y + 0.1), Inches(7.0), Inches(0.4),
                 expr, size=13, color=BODY)
        add_text(s, Inches(11.2), Inches(y + 0.1), Inches(1.5), Inches(0.4),
                 weight, size=13, bold=True, color=ACCENT)

    add_rect(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.05), fill=ACCENT_SOFT)
    add_text(
        s, Inches(0.85), Inches(6.0), Inches(11.6), Inches(0.9),
        "가중치는 PM 피드백 + 플랫폼 similar 신호로 Grid Search 자동 재튜닝.\n"
        "component score 4개를 리포트에 함께 노출 → 왜 경쟁작인지 설명 가능.",
        size=13, color=INK,
    )

    # ---- 6. 솔루션 개요 ----------------------------------------------
    s = new_slide()
    add_title(s, "5개 레이어로 동작합니다", eyebrow="Solution Overview")
    layers = [
        ("① Collectors",  "Steam · App Store · itch.io 메타/리뷰 수집"),
        ("② Pipeline",    "정규화 + Groq LLM feature 추출 + 로컬 임베딩 (384d)"),
        ("③ Engine",      "4축 유사도 · 가중치 튜닝 · weak label 수집"),
        ("④ Report",      "Jinja2 + LLM 요약으로 주간 Markdown 리포트"),
        ("⑤ API + FE",    "FastAPI read-only + Next.js 16 3페이지 (upvote/downvote)"),
    ]
    top = 2.1
    box_h = 0.85
    for i, (name, desc) in enumerate(layers):
        y = top + i * (box_h + 0.1)
        add_rect(s, Inches(0.6), Inches(y), Inches(2.8), Inches(box_h),
                 fill=ACCENT)
        add_text(s, Inches(0.8), Inches(y + 0.22), Inches(2.6), Inches(0.5),
                 name, size=16, bold=True, color=BG)
        add_rect(s, Inches(3.45), Inches(y), Inches(9.25), Inches(box_h),
                 fill=RGBColor(0xF8, 0xFA, 0xFC), line=LINE)
        add_text(s, Inches(3.65), Inches(y + 0.22), Inches(9), Inches(0.5),
                 desc, size=14, color=BODY)

    # ---- 7. PM 사용 플로우 -------------------------------------------
    s = new_slide()
    add_title(s, "PM 은 이렇게 씁니다 (5단계)", eyebrow="User Flow")
    steps = [
        ("1", "내 게임 등록",
         "FE 에서 Steam AppID 입력 → DB 에 is_my_game=TRUE 플래그"),
        ("2", "주간 배치 실행",
         "파이프라인이 내 게임 base 로 Top N 유사도 + 리포트 생성"),
        ("3", "내 게임 목록 확인",
         "홈 '/' 에서 내 게임 카드 그리드 + 각 카드 Competitors/Report 링크"),
        ("4", "경쟁작 분석",
         "'/games/<id>' 에서 Top 10 + component bar + 👍/👎 + 스토어 링크"),
        ("5", "주간 리포트",
         "'/games/<id>/report' 에서 표·신규진입·순위변동·업데이트 요약 + PDF 저장"),
    ]
    top = 2.1
    row_h = 0.95
    for i, (num, title, desc) in enumerate(steps):
        y = top + i * row_h
        add_rect(s, Inches(0.6), Inches(y), Inches(0.8), Inches(0.8),
                 fill=ACCENT)
        add_text(s, Inches(0.6), Inches(y + 0.2), Inches(0.8), Inches(0.5),
                 num, size=22, bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.6), Inches(y + 0.05), Inches(11), Inches(0.4),
                 title, size=16, bold=True, color=INK)
        add_text(s, Inches(1.6), Inches(y + 0.42), Inches(11), Inches(0.4),
                 desc, size=13, color=BODY)

    # ---- 8. 데모 스크린샷 placeholders --------------------------------
    for title, subtitle in [
        ("데모 ① 내 게임 등록 화면", "Home ('/') — 등록 폼 + 내 게임 카드 그리드"),
        ("데모 ② 경쟁작 분석 화면", "/games/<id> — Top 10 + component bar + 스토어 링크"),
        ("데모 ③ 주간 리포트 화면", "/games/<id>/report — 표·신규진입·요약 + PDF 다운로드"),
    ]:
        s = new_slide()
        add_title(s, title, eyebrow="Demo")
        add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.5),
                 subtitle, size=14, color=MUTE)
        add_rect(s, Inches(0.6), Inches(2.7), Inches(12.1), Inches(4.2),
                 fill=RGBColor(0xF4, 0xF4, 0xF5), line=LINE)
        add_text(
            s, Inches(0.6), Inches(4.6), Inches(12.1), Inches(0.4),
            "[ 스크린샷 자리 ]",
            size=18, color=MUTE, align=PP_ALIGN.CENTER,
        )
        add_text(
            s, Inches(0.6), Inches(5.05), Inches(12.1), Inches(0.4),
            "발표 전 교체 — localhost:3000 에서 캡처",
            size=11, color=MUTE, align=PP_ALIGN.CENTER,
        )

    # ---- 9. 피드백 루프 -----------------------------------------------
    s = new_slide()
    add_title(s, "피드백이 다음 주 품질로 이어집니다", eyebrow="Feedback Loop")
    add_text(
        s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.5),
        "PM 에게 추가 업무 없음 — 👍/👎 만 남기면 됩니다.",
        size=16, bold=True, color=INK,
    )
    loop = [
        ("리포트에서 👍/👎",        "FE 클릭 이벤트"),
        ("pm_feedback 테이블",      "DB 로그"),
        ("weight_tuner (주간)",     "NDCG@k grid search"),
        ("weight_history",          "가중치 이력"),
        ("다음 주 유사도 재계산",   "리포트 품질 개선"),
    ]
    y0 = 3.0
    box_w = 2.35
    gap = 0.1
    for i, (top_text, bot_text) in enumerate(loop):
        x = 0.6 + i * (box_w + gap)
        add_rect(s, Inches(x), Inches(y0), Inches(box_w), Inches(2.0),
                 fill=ACCENT_SOFT, line=LINE)
        add_text(s, Inches(x + 0.1), Inches(y0 + 0.3), Inches(box_w - 0.2), Inches(0.9),
                 top_text, size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_text(s, Inches(x + 0.1), Inches(y0 + 1.2), Inches(box_w - 0.2), Inches(0.7),
                 bot_text, size=11, color=MUTE, align=PP_ALIGN.CENTER)
        if i < len(loop) - 1:
            ax = x + box_w + 0.005
            add_text(s, Inches(ax - 0.05), Inches(y0 + 0.85), Inches(0.2), Inches(0.4),
                     "▶", size=14, color=ACCENT, align=PP_ALIGN.CENTER)

    add_text(
        s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.2),
        "HITL '검수·승인' 플로우는 의도적으로 제거.\n"
        "검수를 '업무' 가 아닌 '데이터 생산' 으로 재정의 — PM 의 월요일을 되돌려줌.",
        size=14, color=BODY,
    )

    # ---- 10. 기술 스택 / 아키텍처 요약 --------------------------------
    s = new_slide()
    add_title(s, "기술 스택 한 장", eyebrow="Under the Hood")
    stack = [
        ("언어/런타임", "Python 3.11+ · Node 20 · TypeScript"),
        ("DB",         "PostgreSQL 16 + pgvector (local Docker / 프로덕션 Supabase)"),
        ("LLM",        "Groq llama-3.3-70b-versatile (OpenAI 호환 · 무료 tier)"),
        ("임베딩",     "sentence-transformers all-MiniLM-L6-v2 (384d · 완전 로컬)"),
        ("API",        "FastAPI + uvicorn (read-only + POST /feedback, POST /games/my)"),
        ("FE",         "Next.js 16 + React 19 + Tailwind v4 + turbopack"),
        ("IaC",        "Terraform (Supabase + Vercel 모듈)"),
    ]
    top = 2.1
    row_h = 0.55
    for i, (k, v) in enumerate(stack):
        y = top + i * row_h
        if i % 2 == 1:
            add_rect(s, Inches(0.6), Inches(y), Inches(12.1), Inches(row_h),
                     fill=RGBColor(0xF8, 0xFA, 0xFC))
        add_text(s, Inches(0.85), Inches(y + 0.12), Inches(3.2), Inches(0.4),
                 k, size=13, bold=True, color=INK)
        add_text(s, Inches(4.1), Inches(y + 0.12), Inches(8.5), Inches(0.4),
                 v, size=13, color=BODY)

    # ---- 11. 비용 / 투자 대비 효과 -----------------------------------
    s = new_slide()
    add_title(s, "운영 비용은 사실상 0원입니다", eyebrow="Cost")
    cols = [
        ("LLM",        "Groq 무료 tier",      "$0 / 월"),
        ("임베딩",     "로컬 sentence-transformers", "$0 / 월"),
        ("DB",         "Supabase 무료 프로젝트 (or self-host)", "$0 / 월"),
        ("호스팅",     "Vercel Hobby (FE+API)", "$0 / 월"),
        ("배치 러너",  "로컬 또는 GitHub Actions cron", "$0 / 월"),
    ]
    top = 2.2
    for i, (a, b, c) in enumerate(cols):
        y = top + i * 0.6
        if i % 2 == 1:
            add_rect(s, Inches(0.6), Inches(y), Inches(12.1), Inches(0.6),
                     fill=RGBColor(0xF8, 0xFA, 0xFC))
        add_text(s, Inches(0.85), Inches(y + 0.15), Inches(2.5), Inches(0.4),
                 a, size=14, bold=True, color=INK)
        add_text(s, Inches(3.6), Inches(y + 0.15), Inches(7), Inches(0.4),
                 b, size=13, color=BODY)
        add_text(s, Inches(10.8), Inches(y + 0.15), Inches(2), Inches(0.4),
                 c, size=14, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)

    add_rect(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.3), fill=ACCENT_SOFT)
    add_text(
        s, Inches(0.85), Inches(5.85), Inches(11.5), Inches(0.4),
        "ROI 관점",
        size=13, bold=True, color=ACCENT,
    )
    add_text(
        s, Inches(0.85), Inches(6.25), Inches(11.5), Inches(0.8),
        "PM 1명 × 주 X시간 수동 리서치 → 주 10분으로 단축.\n"
        "도입 비용: 초기 개발 공수 외 운영비 0원 · 확장 시에도 트래픽 비례 선형 증가만.",
        size=13, color=INK,
    )

    # ---- 12. 현재 범위 / 한계 ----------------------------------------
    s = new_slide()
    add_title(s, "현재 범위와 한계", eyebrow="Scope & Limitations")
    add_text(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(0.4),
             "포함됨", size=14, bold=True, color=RGBColor(0x06, 0x5F, 0x46))
    add_bullets(
        s, Inches(0.6), Inches(2.45), Inches(5.8), Inches(4.5),
        [
            "Steam 메타/리뷰 수집 + 내 게임 등록",
            "LLM feature 추출 + 384d 임베딩",
            "4축 유사도 + 주간 스냅샷",
            "Jinja2 Markdown 리포트 + LLM 요약",
            "FE 3페이지 (홈/경쟁작/리포트 + PDF 저장)",
            "PM 피드백 수집 + 가중치 튜닝",
        ],
        size=13,
    )
    add_text(s, Inches(6.9), Inches(2.0), Inches(5.8), Inches(0.4),
             "범위 밖 / 후속", size=14, bold=True, color=RGBColor(0x99, 0x1B, 0x1B))
    add_bullets(
        s, Inches(6.9), Inches(2.45), Inches(5.8), Inches(4.5),
        [
            "인증·권한 (현재 비공개 URL + Vercel password protection 전제)",
            "App Store / Play Store 'add-my-game' 지원",
            "배치 자동화 (GitHub Actions cron)",
            "POST /reports/generate 의 serverless 지원 (ML deps 50MB 제약)",
            "target user overlap 직접 측정 (YouTube / Twitch 신호)",
            "실제 PM 대상 사용자 테스트 + KPI 측정",
        ],
        size=13,
    )

    # ---- 13. 로드맵 (30/60/90) ---------------------------------------
    s = new_slide()
    add_title(s, "로드맵", eyebrow="Next 30 / 60 / 90 Days")
    phases = [
        ("30일",
         RGBColor(0xDB, 0xEA, 0xFE),
         RGBColor(0x1E, 0x3A, 0x8A),
         [
             "실제 PM 2~3명 onboarding",
             "시간 단축 KPI 측정 시작",
             "GitHub Actions 주간 cron",
             "기본 인증 (Vercel pw protection)",
         ]),
        ("60일",
         RGBColor(0xFE, 0xF3, 0xC7),
         RGBColor(0x78, 0x35, 0x0F),
         [
             "App Store add-my-game 지원",
             "가중치 튜닝 A/B 측정",
             "피드백 50+ 확보 후 재학습",
             "리포트 UX 개선 (PM 피드백 반영)",
         ]),
        ("90일",
         RGBColor(0xD1, 0xFA, 0xE5),
         RGBColor(0x06, 0x5F, 0x46),
         [
             "YouTube/Twitch target-user 신호 통합",
             "Play Store 공식 경로 합법성 검토",
             "SSO 연동 + 접근권한 관리",
             "운영 전환 제안서 (비용/SLA/책임)",
         ]),
    ]
    top = 2.1
    box_w = 4.05
    gap = 0.1
    for i, (phase, bg, fg, items) in enumerate(phases):
        x = 0.6 + i * (box_w + gap)
        add_rect(s, Inches(x), Inches(top), Inches(box_w), Inches(0.7), fill=bg)
        add_text(s, Inches(x + 0.2), Inches(top + 0.15), Inches(box_w - 0.4), Inches(0.5),
                 phase, size=18, bold=True, color=fg)
        add_rect(s, Inches(x), Inches(top + 0.7), Inches(box_w), Inches(4.1),
                 fill=RGBColor(0xFA, 0xFA, 0xFA), line=LINE)
        add_bullets(
            s, Inches(x + 0.2), Inches(top + 0.9), Inches(box_w - 0.4), Inches(3.9),
            items, size=13,
        )

    # ---- 14. 의사결정 요청 (Ask) --------------------------------------
    s = new_slide()
    add_title(s, "필요한 의사결정", eyebrow="Ask")
    add_bullets(
        s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(4.0),
        [
            "PoC → 내부 운영 전환 승인 — 초기 PM 3명, 4주 파일럿",
            "외부 신호 확장 검토 — YouTube Data API 키 발급 승인",
            "운영 환경 확정 — Supabase 유료 tier vs 사내 Postgres",
            "성공 기준 합의 — '리포트 준비 시간 X → 1시간' 의 X 실측치 확보 주체",
            "스폰서 지정 — 리서치팀 vs 툴팀 중 어느 쪽이 프로덕트 오너십을 가질지",
        ],
        size=16,
    )
    add_rect(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.0), fill=ACCENT_SOFT)
    add_text(
        s, Inches(0.85), Inches(6.2), Inches(11.5), Inches(0.7),
        "오늘 결정되면 5월 둘째 주부터 실제 PM 대상 파일럿 착수 가능합니다.",
        size=14, bold=True, color=ACCENT,
    )

    # ---- 15. Q&A -----------------------------------------------------
    s = new_slide()
    add_rect(s, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(SLIDE_H_IN),
             fill=ACCENT)
    add_text(
        s, Inches(0.6), Inches(3.0), Inches(12), Inches(1.2),
        "Q & A",
        size=72, bold=True, color=BG, align=PP_ALIGN.CENTER,
    )
    add_text(
        s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.5),
        "질문 · 의견 · 반론 환영합니다",
        size=18, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER,
    )

    # ---- 16. Appendix — DB 스키마 -------------------------------------
    s = new_slide()
    add_title(s, "부록 · DB 스키마 13개 테이블", eyebrow="Appendix")
    tables = [
        ("raw_games",                 "플랫폼 원본 payload (JSONB)"),
        ("raw_reviews",               "리뷰 원본 (text, rating, posted_at)"),
        ("games",                     "정규화 게임 · is_my_game 플래그"),
        ("game_features",             "LLM 추출 feature · SCD Type 2"),
        ("game_embeddings",           "pgvector vector(384)"),
        ("game_similarities_weekly",  "주간 유사도 스냅샷 + component_scores"),
        ("pm_feedback",               "upvote/downvote/clicked/added"),
        ("llm_cache / embedding_cache", "호출 캐시 (SHA256 키)"),
        ("pipeline_runs",             "stage × week × status 관측성"),
        ("weekly_reports",            "생성된 리포트 (JSONB)"),
        ("weight_history",            "튜닝된 가중치 + NDCG@10"),
        ("weak_similarities",         "플랫폼 similar signal (tag overlap 등)"),
    ]
    top = 2.05
    row_h = 0.4
    for i, (name, desc) in enumerate(tables):
        col = i % 2
        row = i // 2
        x = 0.6 + col * 6.15
        y = top + row * row_h
        add_text(s, Inches(x), Inches(y), Inches(2.4), Inches(0.35),
                 name, size=12, bold=True, color=ACCENT, font="Consolas")
        add_text(s, Inches(x + 2.5), Inches(y), Inches(3.6), Inches(0.35),
                 desc, size=12, color=BODY)

    # ---- 17. Appendix — 유사도 공식 -----------------------------------
    s = new_slide()
    add_title(s, "부록 · 유사도 공식", eyebrow="Appendix")
    add_rect(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.4),
             fill=RGBColor(0x0F, 0x17, 0x2A))
    add_text(
        s, Inches(0.85), Inches(2.4), Inches(11.5), Inches(1.0),
        "score = 0.40 · semantic + 0.25 · genre + 0.20 · tier + 0.15 · bm",
        size=22, bold=True, color=BG, font="Consolas",
    )
    axes = [
        ("semantic",  "cosine(desc_embed_a, desc_embed_b)",
         "description + reviews 를 합쳐 임베딩한 벡터의 코사인 유사도"),
        ("genre",     "cosine(genre_embed_a, genre_embed_b)",
         "LLM 이 추출한 장르 문자열의 임베딩 — 카테고리 이진 매칭이 아닌 연속값"),
        ("tier",      "exp(-|tier_a - tier_b| / 0.3)",
         "log1p(review_count) 정규화 후 거리 기반 Gaussian — 체급 유사도"),
        ("bm",        "1 - KL_sym(bm_a, bm_b) / (2·log 4)",
         "gacha/ads/premium/sub 확률분포의 대칭 KL 거리 → 유사도로 반전"),
    ]
    top = 3.8
    for i, (name, formula, desc) in enumerate(axes):
        y = top + i * 0.7
        add_text(s, Inches(0.6), Inches(y), Inches(1.6), Inches(0.3),
                 name, size=13, bold=True, color=ACCENT, font="Consolas")
        add_text(s, Inches(2.3), Inches(y), Inches(5.0), Inches(0.3),
                 formula, size=11, color=INK, font="Consolas")
        add_text(s, Inches(7.4), Inches(y), Inches(5.3), Inches(0.6),
                 desc, size=11, color=BODY)

    return slides


# ------------------------------------------------------------------
# Main
# ------------------------------------------------------------------

def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    slides = build(prs, page_total_provider=None)

    # Pass 2: footers now that we know total
    total = len(slides)
    for i, s in enumerate(slides, start=1):
        # Skip cover (1) and Q&A (which is the full-bleed accent slide).
        # Footer added to all content slides except cover, Q&A, and appendix cover-style.
        if i == 1:
            continue
        # Q&A is the solid accent one — detect by checking if index matches ours.
        # We placed Q&A at slide 15. We'll just skip by checking bg via heuristic:
        # Simpler: add footer to all except slide 1 and 15.
        if i == 15:
            continue
        add_footer(s, i, total)

    out_dir = Path(__file__).resolve().parents[1] / "docs"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "presentation.pptx"
    prs.save(out_path)
    print(f"✓ wrote {out_path} ({total} slides)")


if __name__ == "__main__":
    main()
