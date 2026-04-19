"""Build 과제전형 submission deck.

6 required sections (① 문제 정의 · ② 해결 방법 · ③ AI 도구 · ④ 프로토타이핑 계획
· ⑤ 프로토타이핑 방법 · ⑥ 바이브코딩 지시) + 3 wireframes + prototype/demo slides.

Run:
    python scripts/build_submission_deck.py

Output:
    docs/submission/submission.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ------------------------------------------------------------------
# Submission metadata (update here, then rebuild)
# ------------------------------------------------------------------
APPLICANT_NAME = "윤동혁"
APPLICANT_POSITION = "AI 애플리케이션 엔지니어"
SUBMISSION_DATE = "2026-04-21"

# Replace after `terraform apply` and demo recording.
PROTOTYPE_WEB_URL = "https://gca-web.vercel.app"      # placeholder
PROTOTYPE_API_URL = "https://gca-api.vercel.app"      # placeholder
GITHUB_REPO_URL = "https://github.com/hyukyyy/game-competitor-analysis-poc"
DEMO_GIF_URL = "docs/submission/demo.gif"             # user-recorded local file

# ------------------------------------------------------------------
# Palette (mirrors build_deck.py for consistency)
# ------------------------------------------------------------------
INK = RGBColor(0x18, 0x18, 0x1B)
BODY = RGBColor(0x3F, 0x3F, 0x46)
MUTE = RGBColor(0x71, 0x71, 0x7A)
LINE = RGBColor(0xE4, 0xE4, 0xE7)
BG = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x1F, 0x2A, 0x44)
ACCENT_SOFT = RGBColor(0xEE, 0xF2, 0xF7)
SURFACE = RGBColor(0xF8, 0xFA, 0xFC)
WIRE_FILL = RGBColor(0xF4, 0xF4, 0xF5)

FONT = "Calibri"
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------

def set_bg(slide, rgb=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=BODY, align=PP_ALIGN.LEFT, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=16, color=BODY, bullet="•", line_spacing=1.35):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_hline(slide, left, top, width, color=LINE, weight=0.75):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_rect(slide, left, top, width, height, *,
             fill=ACCENT_SOFT, line=None, shape=MSO_SHAPE.RECTANGLE):
    rect = slide.shapes.add_shape(shape, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(0.5)
    rect.shadow.inherit = False
    return rect


def add_footer(slide, page_num, total):
    add_text(slide,
             Inches(0.5), Inches(SLIDE_H_IN - 0.35), Inches(10), Inches(0.25),
             f"{APPLICANT_NAME} · {APPLICANT_POSITION} · 과제전형 제출 · {SUBMISSION_DATE}",
             size=9, color=MUTE)
    add_text(slide,
             Inches(SLIDE_W_IN - 1.5), Inches(SLIDE_H_IN - 0.35), Inches(1.0), Inches(0.25),
             f"{page_num} / {total}",
             size=9, color=MUTE, align=PP_ALIGN.RIGHT)


def add_title(slide, text, *, eyebrow=None):
    top = Inches(0.55)
    if eyebrow:
        add_text(slide, Inches(0.6), top, Inches(12), Inches(0.3),
                 eyebrow.upper(), size=10, bold=True, color=ACCENT)
        top = Inches(0.9)
    add_text(slide, Inches(0.6), top, Inches(12.1), Inches(0.7),
             text, size=30, bold=True, color=INK)
    add_hline(slide, Inches(0.6), Inches(1.75), Inches(12.1))


def add_section_badge(slide, left, top, label):
    """Pill-shaped section marker like '① 문제 정의'."""
    add_rect(slide, left, top, Inches(2.2), Inches(0.4),
             fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, left, top + Emu(50000), Inches(2.2), Inches(0.35),
             label, size=12, bold=True, color=BG, align=PP_ALIGN.CENTER)


# ------------------------------------------------------------------
# Wireframe helpers
# ------------------------------------------------------------------

def _wire_frame(slide, left, top, width, height):
    """Browser-like outer frame for wireframes."""
    add_rect(slide, left, top, width, height, fill=WIRE_FILL, line=LINE)
    # top bar (window chrome)
    bar_h = Inches(0.3)
    add_rect(slide, left, top, width, bar_h,
             fill=RGBColor(0xE4, 0xE4, 0xE7))


def _wire_box(slide, left, top, width, height, label, *,
              sub=None, fill=BG, line_color=RGBColor(0xA1, 0xA1, 0xAA)):
    add_rect(slide, left, top, width, height, fill=fill, line=line_color)
    add_text(slide, left + Inches(0.1), top + Inches(0.08), width - Inches(0.2), Inches(0.3),
             label, size=11, bold=True, color=INK)
    if sub:
        add_text(slide, left + Inches(0.1), top + Inches(0.38), width - Inches(0.2), Inches(0.5),
                 sub, size=9, color=MUTE)


def _draw_wireframe_home(slide, left, top, width, height):
    _wire_frame(slide, left, top, width, height)
    # URL bar
    add_text(slide, left + Inches(0.1), top + Inches(0.06), Inches(3), Inches(0.2),
             "gca-web.vercel.app /", size=9, color=MUTE, font="Consolas")

    inner_top = top + Inches(0.45)
    # Header
    _wire_box(slide, left + Inches(0.2), inner_top, width - Inches(0.4), Inches(0.5),
              "🎮 Game Competitor Analysis",
              sub="Home · Competitors · Reports")

    # AddMyGameForm card
    form_top = inner_top + Inches(0.65)
    _wire_box(slide, left + Inches(0.2), form_top, width - Inches(0.4), Inches(1.0),
              "＋ Add My Game",
              sub="[ Steam AppID input ___________ ]    [ Submit ]",
              fill=ACCENT_SOFT)

    # My Games grid (2x3)
    grid_top = form_top + Inches(1.15)
    card_w = (width - Inches(0.6)) / 3
    card_h = Inches(1.0)
    for row in range(2):
        for col in range(3):
            gx = left + Inches(0.2) + col * (card_w + Inches(0.1))
            gy = grid_top + row * (card_h + Inches(0.1))
            _wire_box(slide, gx, gy, card_w, card_h,
                      f"Game card {row * 3 + col + 1}",
                      sub="title · genre · tier · is_my_game")


def _draw_wireframe_competitors(slide, left, top, width, height):
    _wire_frame(slide, left, top, width, height)
    add_text(slide, left + Inches(0.1), top + Inches(0.06), Inches(4), Inches(0.2),
             "gca-web.vercel.app / games / [id]", size=9, color=MUTE, font="Consolas")

    inner_top = top + Inches(0.45)
    # Base game banner
    _wire_box(slide, left + Inches(0.2), inner_top, width - Inches(3.0), Inches(0.75),
              "Base game · 내 게임 (is_my_game)",
              sub="title · store link · generated @ 2026-MM-DD")
    # Weekly Report CTA
    _wire_box(slide, left + width - Inches(2.7), inner_top, Inches(2.5), Inches(0.75),
              "→ Weekly Report",
              sub="/games/[id]/report",
              fill=ACCENT, line_color=ACCENT)

    # Competitor table header
    table_top = inner_top + Inches(0.9)
    _wire_box(slide, left + Inches(0.2), table_top, width - Inches(0.4), Inches(0.35),
              "#  Competitor  |  Score  |  semantic • genre • tier • BM  |  Vote",
              fill=ACCENT_SOFT)

    # Table rows with component bar + upvote
    row_h = Inches(0.45)
    for i in range(6):
        ry = table_top + Inches(0.4) + i * (row_h + Inches(0.05))
        # row background
        _wire_box(slide, left + Inches(0.2), ry, width - Inches(0.4), row_h,
                  f"{i+1}. Competitor #{i+1}",
                  sub="0.8X  ▮▮▮▮▯▯  semantic .40 / genre .25 / tier .20 / BM .15   👍 👎")


def _draw_wireframe_report(slide, left, top, width, height):
    _wire_frame(slide, left, top, width, height)
    add_text(slide, left + Inches(0.1), top + Inches(0.06), Inches(4.5), Inches(0.2),
             "gca-web.vercel.app / games / [id] / report", size=9, color=MUTE, font="Consolas")

    inner_top = top + Inches(0.45)
    # Header + Download PDF button
    _wire_box(slide, left + Inches(0.2), inner_top, width - Inches(2.7), Inches(0.7),
              "Weekly Competitor Report",
              sub="base: 내 게임  ·  generated: 2026-MM-DD  ·  Top 10 competitors")
    _wire_box(slide, left + width - Inches(2.4), inner_top, Inches(2.2), Inches(0.7),
              "⬇ Download PDF",
              sub="window.print() · @media print hides nav",
              fill=ACCENT, line_color=ACCENT)

    # Markdown sections
    section_top = inner_top + Inches(0.85)
    sections = [
        ("## Top Movers",        "순위 상승 · 신규 진입 경쟁작 리스트 (표)"),
        ("## BM Analysis",       "gacha / ads / premium / sub 분포 요약"),
        ("## Feature Trends",    "LLM 추출 feature 의 주간 변화"),
        ("## Recommendations",   "PM 을 위한 다음 주 액션 아이템"),
    ]
    sh = Inches(0.55)
    for i, (title, desc) in enumerate(sections):
        sy = section_top + i * (sh + Inches(0.08))
        _wire_box(slide, left + Inches(0.2), sy, width - Inches(0.4), sh,
                  title, sub=desc)


# ------------------------------------------------------------------
# Slide builders
# ------------------------------------------------------------------

def build(prs):
    slides = []

    def new_slide():
        layout = prs.slide_layouts[6]
        s = prs.slides.add_slide(layout)
        set_bg(s)
        slides.append(s)
        return s

    # -------------------------------------------------- 1. Cover
    s = new_slide()
    add_rect(s, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(SLIDE_H_IN), fill=ACCENT)
    add_text(s, Inches(0.8), Inches(0.8), Inches(11), Inches(0.4),
             "과제전형 제출",
             size=12, bold=True, color=RGBColor(0xCB, 0xD5, 0xE1))
    add_text(s, Inches(0.8), Inches(1.3), Inches(12), Inches(1.4),
             "Game Competitor Analysis",
             size=48, bold=True, color=BG)
    add_text(s, Inches(0.8), Inches(2.4), Inches(12), Inches(0.6),
             "AI 기반 주간 경쟁작 리서치 자동화",
             size=22, color=RGBColor(0xE4, 0xE8, 0xF0))

    add_rect(s, Inches(0.8), Inches(4.2), Inches(8), Inches(2.0),
             fill=RGBColor(0x2D, 0x3A, 0x5A))
    add_text(s, Inches(1.0), Inches(4.4), Inches(7.5), Inches(0.4),
             "지원자", size=11, bold=True, color=RGBColor(0xCB, 0xD5, 0xE1))
    add_text(s, Inches(1.0), Inches(4.75), Inches(7.5), Inches(0.5),
             APPLICANT_NAME, size=24, bold=True, color=BG)
    add_text(s, Inches(1.0), Inches(5.35), Inches(7.5), Inches(0.4),
             "지원 포지션", size=11, bold=True, color=RGBColor(0xCB, 0xD5, 0xE1))
    add_text(s, Inches(1.0), Inches(5.7), Inches(7.5), Inches(0.5),
             APPLICANT_POSITION, size=18, color=BG)

    add_text(s, Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
             f"제출일 · {SUBMISSION_DATE}",
             size=12, color=RGBColor(0xCB, 0xD5, 0xE1))

    # -------------------------------------------------- 2. 한 장 요약
    s = new_slide()
    add_title(s, "한 장 요약", eyebrow="Executive Summary")
    add_rect(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.4), fill=ACCENT_SOFT)
    add_text(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.0),
             "게임 PM 이 매주 월요일 수시간씩 수작업으로 하던 경쟁작 리서치를,\n"
             "4축 유사도 엔진 + 주간 배치 + 피드백 루프로 자동화한 AI 제품 PoC.",
             size=18, bold=True, color=INK)
    add_bullets(s, Inches(0.6), Inches(3.9), Inches(12.1), Inches(3.2),
                [
                    "해결 대상: 게임 PM/기획자 · 주간 경쟁작 모니터링 업무",
                    "AI 도구: Groq llama-3.3-70b (feature 추출·요약) + sentence-transformers (384d) + pgvector",
                    "구현 상태: 13 테이블 스키마, 15개 CLI 서브커맨드, 3개 FE 라우트, Terraform IaC",
                    "검증: 5단계 User Flow · 실제 배포 프로토타입 · 데모 GIF",
                    "바이브코딩: Claude Code(Opus) 에게 역할·기능·UI·기술 스택을 포함한 지시 4종",
                ],
                size=15)

    # -------------------------------------------------- 3. ① 문제 정의
    s = new_slide()
    add_title(s, "해결하고자 하는 문제", eyebrow="Section ①  ·  문제 정의")
    add_section_badge(s, Inches(10.8), Inches(0.6), "① 문제 정의")

    # Left: 현재
    add_rect(s, Inches(0.6), Inches(2.1), Inches(5.9), Inches(4.8),
             fill=RGBColor(0xFE, 0xF2, 0xF2), line=RGBColor(0xFE, 0xCA, 0xCA))
    add_text(s, Inches(0.85), Inches(2.3), Inches(5.5), Inches(0.4),
             "현재 업무 방식", size=12, bold=True, color=RGBColor(0x99, 0x1B, 0x1B))
    add_bullets(s, Inches(0.85), Inches(2.75), Inches(5.5), Inches(4.0),
                [
                    "대상 사용자: 게임 PM · 기획자 · 리서치팀",
                    "매주 월요일 아침 수동 경쟁작 조사",
                    "Steam 스토어 검색 → 장르·BM·리뷰 훑기 → 스프레드시트 정리",
                    "담당 게임 수 × 수 시간 → 월요일이 통째로 녹음",
                    "개인 경험 의존 → 재현성·비교성 낮음",
                    "일회성 메모 → 기관 기억으로 축적 안 됨",
                ],
                size=13, color=RGBColor(0x7F, 0x1D, 0x1D))

    # Right: 원하는
    add_rect(s, Inches(6.8), Inches(2.1), Inches(5.9), Inches(4.8),
             fill=RGBColor(0xEC, 0xFD, 0xF5), line=RGBColor(0xA7, 0xF3, 0xD0))
    add_text(s, Inches(7.05), Inches(2.3), Inches(5.5), Inches(0.4),
             "원하는 상태 (문제 해결 후)", size=12, bold=True, color=RGBColor(0x06, 0x5F, 0x46))
    add_bullets(s, Inches(7.05), Inches(2.75), Inches(5.5), Inches(4.0),
                [
                    "월요일 아침 10분 안에 경쟁 구도 파악",
                    "시스템이 후보 Top N 제시 → PM 이 취사선택",
                    "피드백이 다음 주 품질 개선으로 자동 반영",
                    "주간 스냅샷 DB 누적 → 시장 추이 추적 가능",
                    "담당 게임 늘어도 PM 공수 선형 증가 없음",
                    "북극성 지표: 리포트 준비 시간 X → 1시간 이내",
                ],
                size=13, color=RGBColor(0x06, 0x5F, 0x46))

    # -------------------------------------------------- 4. ② 해결 방법 - 제품 개요
    s = new_slide()
    add_title(s, "AI 기반 제품 개요", eyebrow="Section ②  ·  해결 방법 (1/3)")
    add_section_badge(s, Inches(10.8), Inches(0.6), "② 해결 방법")

    add_rect(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.1), fill=ACCENT_SOFT)
    add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.8),
             "내 게임을 등록하면 주간 배치가 4축 유사도로 Top 10 경쟁작을 자동 산출,\n"
             "LLM 이 Markdown 리포트를 생성, PM 은 👍/👎 만 남기면 다음 주 가중치가 학습됨.",
             size=15, bold=True, color=INK)

    # 4축 유사도 표
    head_top = Inches(3.5)
    add_rect(s, Inches(0.6), head_top, Inches(12.1), Inches(0.5), fill=ACCENT)
    add_text(s, Inches(0.85), head_top + Emu(40000), Inches(3.0), Inches(0.4),
             "유사도 축", size=12, bold=True, color=BG)
    add_text(s, Inches(4.2), head_top + Emu(40000), Inches(7.0), Inches(0.4),
             "표현 방식", size=12, bold=True, color=BG)
    add_text(s, Inches(11.2), head_top + Emu(40000), Inches(1.5), Inches(0.4),
             "가중치 (초기)", size=12, bold=True, color=BG)

    rows = [
        ("Semantic", "description + review 임베딩 cosine", "0.40"),
        ("Genre",    "LLM 추출 장르 임베딩 cosine", "0.25"),
        ("Tier",     "log1p(review_count) Gaussian", "0.20"),
        ("BM",       "{gacha, ads, premium, sub} KL 대칭거리", "0.15"),
    ]
    row_h = 0.5
    for i, (axis, expr, weight) in enumerate(rows):
        y = 4.0 + i * row_h
        if i % 2 == 1:
            add_rect(s, Inches(0.6), Inches(y), Inches(12.1), Inches(row_h), fill=SURFACE)
        add_text(s, Inches(0.85), Inches(y + 0.1), Inches(3.0), Inches(0.4),
                 axis, size=13, bold=True, color=INK)
        add_text(s, Inches(4.2), Inches(y + 0.1), Inches(7.0), Inches(0.4),
                 expr, size=13, color=BODY)
        add_text(s, Inches(11.2), Inches(y + 0.1), Inches(1.5), Inches(0.4),
                 weight, size=13, bold=True, color=ACCENT)

    add_text(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.4),
             "score = 0.40·semantic + 0.25·genre + 0.20·tier + 0.15·bm · 가중치는 PM 피드백으로 주간 재튜닝",
             size=12, color=MUTE, font="Consolas")

    # -------------------------------------------------- 5. ② 해결 방법 - 핵심 기능 3개
    s = new_slide()
    add_title(s, "핵심 기능 3개", eyebrow="Section ②  ·  해결 방법 (2/3)")
    add_section_badge(s, Inches(10.8), Inches(0.6), "② 해결 방법")

    features = [
        ("①",
         "내 게임 등록 + 주간 경쟁작 Top 10",
         "Steam AppID 등록 → 배치 파이프라인이 is_my_game=TRUE 게임만 base 로 유사도 계산 →\n"
         "FE 에서 종합 점수 · 4개 component score bar · 스토어 링크 함께 표시."),
        ("②",
         "LLM 주간 Markdown 리포트 + PDF 저장",
         "Jinja2 템플릿 + Groq llama-3.3-70b 요약으로 '신규 진입 · 순위 변동 · BM 분포' 리포트 자동 생성.\n"
         "FE `/games/[id]/report` 에서 react-markdown 렌더링 + window.print() 기반 PDF 저장."),
        ("③",
         "PM 피드백 루프 (가중치 자동 학습)",
         "👍/👎 클릭 → pm_feedback 테이블 → 주간 weight_tuner 가 NDCG@k grid search →\n"
         "weight_history 에 저장 → 다음 주 유사도 계산에 반영. PM 추가 업무 없음."),
    ]
    top = 2.1
    box_h = 1.55
    for i, (num, title, desc) in enumerate(features):
        y = top + i * (box_h + 0.1)
        add_rect(s, Inches(0.6), Inches(y), Inches(1.2), Inches(box_h), fill=ACCENT)
        add_text(s, Inches(0.6), Inches(y + 0.45), Inches(1.2), Inches(0.6),
                 num, size=36, bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_rect(s, Inches(1.85), Inches(y), Inches(10.85), Inches(box_h),
                 fill=SURFACE, line=LINE)
        add_text(s, Inches(2.05), Inches(y + 0.2), Inches(10.5), Inches(0.5),
                 title, size=17, bold=True, color=INK)
        add_text(s, Inches(2.05), Inches(y + 0.65), Inches(10.5), Inches(0.9),
                 desc, size=12, color=BODY)

    # -------------------------------------------------- 6. ② 해결 방법 - User Flow
    s = new_slide()
    add_title(s, "User Flow (5단계)", eyebrow="Section ②  ·  해결 방법 (3/3)")
    add_section_badge(s, Inches(10.8), Inches(0.6), "② 해결 방법")

    steps = [
        ("1", "내 게임 등록",
         "FE 에서 Steam AppID 입력 → DB 에 is_my_game=TRUE 플래그"),
        ("2", "주간 배치 실행",
         "파이프라인이 내 게임 base 로 Top N 유사도 + 리포트 생성"),
        ("3", "내 게임 목록 확인",
         "홈 '/' 카드 그리드 + 각 카드 → Competitors/Report 링크"),
        ("4", "경쟁작 분석",
         "'/games/[id]' Top 10 + component bar + 👍/👎 + 스토어 링크"),
        ("5", "주간 리포트 소비",
         "'/games/[id]/report' markdown 리포트 + Download PDF"),
    ]
    top = 2.1
    row_h = 0.95
    for i, (num, title, desc) in enumerate(steps):
        y = top + i * row_h
        add_rect(s, Inches(0.6), Inches(y), Inches(0.8), Inches(0.8), fill=ACCENT)
        add_text(s, Inches(0.6), Inches(y + 0.2), Inches(0.8), Inches(0.5),
                 num, size=22, bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.6), Inches(y + 0.05), Inches(11), Inches(0.4),
                 title, size=16, bold=True, color=INK)
        add_text(s, Inches(1.6), Inches(y + 0.42), Inches(11), Inches(0.4),
                 desc, size=13, color=BODY)

    # -------------------------------------------------- 7. ③ AI 도구
    s = new_slide()
    add_title(s, "사용할 AI 도구 3종", eyebrow="Section ③  ·  사용할 AI 도구")
    add_section_badge(s, Inches(10.8), Inches(0.6), "③ AI 도구")

    tools = [
        ("Groq · llama-3.3-70b-versatile",
         "LLM (OpenAI API 호환 · 무료 tier)",
         [
             "게임 feature 추출: description/리뷰 → {genres, bm_mix, session_type, core_loop, ...}",
             "주간 리포트 요약: Jinja2 템플릿 섹션별 LLM 요약 생성",
             "역할: 비구조화 텍스트 → 구조화 JSON + 자연어 요약",
         ]),
        ("sentence-transformers · all-MiniLM-L6-v2",
         "임베딩 모델 (완전 로컬 · 384d)",
         [
             "description + 리뷰 합본을 384차원 벡터로 인코딩",
             "장르 텍스트도 동일 모델로 임베딩 (카테고리 이진 매칭 X, 연속값)",
             "역할: 의미 기반 유사도 계산의 재료 벡터 생산",
         ]),
        ("pgvector · PostgreSQL 확장",
         "벡터 DB (Supabase 내장)",
         [
             "game_embeddings 테이블의 vector(384) 컬럼에 저장",
             "cosine distance 로 Top N 후보 선별 → Engine 이 4축 가중합 재계산",
             "역할: 벡터 인덱스 + 스냅샷 가능한 관계형 DB 의 결합",
         ]),
    ]
    top = 2.1
    box_h = 1.55
    for i, (name, role, items) in enumerate(tools):
        y = top + i * (box_h + 0.1)
        add_rect(s, Inches(0.6), Inches(y), Inches(4.0), Inches(box_h), fill=ACCENT_SOFT)
        add_text(s, Inches(0.8), Inches(y + 0.2), Inches(3.8), Inches(0.5),
                 name, size=14, bold=True, color=ACCENT, font="Consolas")
        add_text(s, Inches(0.8), Inches(y + 0.65), Inches(3.8), Inches(0.8),
                 role, size=11, color=BODY)
        add_rect(s, Inches(4.65), Inches(y), Inches(8.05), Inches(box_h),
                 fill=SURFACE, line=LINE)
        add_bullets(s, Inches(4.85), Inches(y + 0.15), Inches(7.8), Inches(box_h - 0.3),
                    items, size=11, line_spacing=1.25)

    # -------------------------------------------------- 8. ④ 프로토타이핑 계획
    s = new_slide()
    add_title(s, "프로토타이핑 계획", eyebrow="Section ④  ·  프로토타이핑 계획")
    add_section_badge(s, Inches(10.8), Inches(0.6), "④ 프로토타이핑 계획")

    plan = [
        ("목적",
         "PM 의 주간 리포트 준비 시간 단축 가설 검증 — '10분 내 경쟁 구도 파악' 가능한지"),
        ("범위",
         "Steam 스토어 Top ~200 게임 풀. 주간 배치 1회. 인증 없음. 내 게임만 base."),
        ("도구 스택",
         "Python CLI · FastAPI · Next.js 16 · PostgreSQL 16 + pgvector · Groq · Vercel + Supabase"),
        ("입력",
         "내 게임 Steam AppID (예: 730=CS2, 1063730=New World)"),
        ("출력",
         "경쟁작 Top 10 (4축 component score 포함) + 주간 Markdown 리포트 + PDF"),
        ("성공 지표",
         "북극성: 리포트 준비 시간 X → 1시간. 보조: upvote rate ≥ 50%, feature 자기일관성 ≥ 90%"),
        ("기간",
         "4주 PoC — Week 1 스키마+수집, Week 2 feature+임베딩, Week 3 유사도+리포트, Week 4 FE+배포"),
    ]
    top = 2.1
    row_h = 0.65
    for i, (k, v) in enumerate(plan):
        y = top + i * row_h
        if i % 2 == 1:
            add_rect(s, Inches(0.6), Inches(y), Inches(12.1), Inches(row_h), fill=SURFACE)
        add_text(s, Inches(0.85), Inches(y + 0.17), Inches(2.0), Inches(0.4),
                 k, size=13, bold=True, color=ACCENT)
        add_text(s, Inches(2.95), Inches(y + 0.17), Inches(9.7), Inches(0.4),
                 v, size=12, color=BODY)

    # -------------------------------------------------- 9. ⑤ 방법 - 아키텍처 5레이어
    s = new_slide()
    add_title(s, "5 레이어 아키텍처", eyebrow="Section ⑤  ·  프로토타이핑 방법 (1/3)")
    add_section_badge(s, Inches(10.8), Inches(0.6), "⑤ 방법")

    layers = [
        ("① Collectors", "Steam (Store API + Reviews API) 수집 + appstore/itch stub"),
        ("② Pipeline",   "normalize → Groq feature 추출 → sentence-transformers 임베딩 (384d)"),
        ("③ Engine",     "4축 유사도 계산 · weak label 수집 · weight_tuner (NDCG@k)"),
        ("④ Report",     "Jinja2 템플릿 + Groq LLM 섹션별 요약 → weekly_reports 저장"),
        ("⑤ API + FE",   "FastAPI read-only + POST /feedback · Next.js 16 App Router 3페이지"),
    ]
    top = 2.1
    box_h = 0.85
    for i, (name, desc) in enumerate(layers):
        y = top + i * (box_h + 0.1)
        add_rect(s, Inches(0.6), Inches(y), Inches(3.2), Inches(box_h), fill=ACCENT)
        add_text(s, Inches(0.8), Inches(y + 0.22), Inches(3.0), Inches(0.5),
                 name, size=16, bold=True, color=BG)
        add_rect(s, Inches(3.85), Inches(y), Inches(8.85), Inches(box_h),
                 fill=SURFACE, line=LINE)
        add_text(s, Inches(4.05), Inches(y + 0.22), Inches(8.6), Inches(0.5),
                 desc, size=13, color=BODY)

    # -------------------------------------------------- 10. ⑤ DB 스키마
    s = new_slide()
    add_title(s, "DB 스키마 13개 테이블", eyebrow="Section ⑤  ·  프로토타이핑 방법 (2/3)")
    add_section_badge(s, Inches(10.8), Inches(0.6), "⑤ 방법")

    tables = [
        ("raw_games",                 "플랫폼 원본 payload (JSONB)"),
        ("raw_reviews",               "리뷰 원본 (text, rating, posted_at)"),
        ("games",                     "정규화 게임 · is_my_game 플래그"),
        ("game_features",             "LLM 추출 feature · SCD Type 2"),
        ("game_embeddings",           "pgvector vector(384)"),
        ("game_similarities_weekly",  "주간 유사도 + component_scores"),
        ("pm_feedback",               "upvote/downvote/clicked/added"),
        ("llm_cache",                 "LLM 호출 캐시 (SHA256 키)"),
        ("embedding_cache",           "임베딩 캐시 (SHA256 키)"),
        ("pipeline_runs",             "stage × week × status 관측성"),
        ("weekly_reports",            "생성된 리포트 (JSONB)"),
        ("weight_history",            "튜닝된 가중치 + NDCG@10"),
        ("weak_similarities",         "플랫폼 similar (tag overlap 등)"),
    ]
    top = 2.05
    row_h = 0.38
    for i, (name, desc) in enumerate(tables):
        col = i % 2
        row = i // 2
        x = 0.6 + col * 6.15
        y = top + row * row_h
        add_text(s, Inches(x), Inches(y), Inches(2.6), Inches(0.35),
                 name, size=12, bold=True, color=ACCENT, font="Consolas")
        add_text(s, Inches(x + 2.65), Inches(y), Inches(3.5), Inches(0.35),
                 desc, size=11, color=BODY)

    # -------------------------------------------------- 11. ⑤ 주간 배치 파이프라인
    s = new_slide()
    add_title(s, "주간 배치 파이프라인", eyebrow="Section ⑤  ·  프로토타이핑 방법 (3/3)")
    add_section_badge(s, Inches(10.8), Inches(0.6), "⑤ 방법")

    pipe = [
        ("collect:steam",       "Steam Store/Reviews 수집 → raw_games, raw_reviews"),
        ("normalize",           "플랫폼 스키마 차이 정규화 → games"),
        ("extract-features",    "Groq LLM feature 추출 → game_features (SCD Type 2)"),
        ("embed",               "sentence-transformers 384d → game_embeddings"),
        ("weak-labels",         "tag overlap 등 약한 신호 → weak_similarities"),
        ("tune-weights",        "pm_feedback × weak_similarities → NDCG grid search"),
        ("similarity --week",   "내 게임 base 4축 유사도 → game_similarities_weekly"),
        ("report --week",       "Jinja2 + Groq 요약 → weekly_reports"),
    ]
    top = 2.1
    row_h = 0.55
    for i, (cmd, desc) in enumerate(pipe):
        y = top + i * row_h
        if i % 2 == 1:
            add_rect(s, Inches(0.6), Inches(y), Inches(12.1), Inches(row_h), fill=SURFACE)
        add_text(s, Inches(0.85), Inches(y + 0.13), Inches(3.5), Inches(0.4),
                 f"gca {cmd}", size=12, bold=True, color=ACCENT, font="Consolas")
        add_text(s, Inches(4.45), Inches(y + 0.13), Inches(8.2), Inches(0.4),
                 desc, size=12, color=BODY)

    add_rect(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.55), fill=ACCENT_SOFT)
    add_text(s, Inches(0.85), Inches(6.8), Inches(11.5), Inches(0.4),
             "배치 시점: 매주 월요일 오전 (GitHub Actions cron 또는 로컬). "
             "실행 시간: ~5분 (캐시 적중 시). 실패 시 pipeline_runs 에 상태 기록.",
             size=11, color=INK)

    # -------------------------------------------------- 12. ⑥ 바이브코딩 지시 (1/2)
    s = new_slide()
    add_title(s, "바이브코딩 지시 — 스캐폴딩 · FE", eyebrow="Section ⑥  ·  바이브코딩 지시 (1/2)")
    add_section_badge(s, Inches(10.8), Inches(0.6), "⑥ 바이브코딩")

    add_text(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.35),
             "도구: Claude Code (Opus 4.7) · 각 프롬프트에 [역할 · 기능 · UI/UX · 기술 스택] 4요소 포함",
             size=11, color=MUTE)

    prompt1 = (
        "역할: 백엔드 아키텍트.\n"
        "기능: Steam/AppStore/itch 에서 메타+리뷰를 수집해 4축(semantic/genre/tier/BM) 유사도로 "
        "내 게임의 경쟁작 Top N 을 주간 배치로 계산·저장.\n"
        "기술 스택: Python 3.11 + FastAPI + PostgreSQL 16(pgvector) + Groq(OpenAI 호환) + "
        "sentence-transformers all-MiniLM-L6-v2.\n"
        "산출물: ① 13개 테이블 스키마(schema.sql) ② collectors → normalize → feature_extractor(LLM) "
        "→ embedder → similarity 파이프라인 ③ gca CLI 서브커맨드 ④ llm_cache/embedding_cache 로 비용 제어."
    )
    prompt2 = (
        "역할: 프런트엔드 엔지니어.\n"
        "기능: 내 게임 등록 · 경쟁작 Top 10 표시 · 👍/👎 피드백 · 주간 리포트 소비.\n"
        "UI/UX: 3개 라우트 — `/`(내 게임 카드 그리드 + AddMyGameForm), "
        "`/games/[id]`(경쟁작 테이블 + 4축 component score bar + upvote/downvote), "
        "`/games/[id]/report`(markdown 리포트 + 우상단 Download PDF 버튼). "
        "톤: Tailwind zinc 모노톤 · 서버 컴포넌트 우선.\n"
        "기술 스택: Next.js 16 App Router + React 19 + Tailwind v4 + turbopack · react-markdown + remark-gfm · window.print() PDF."
    )

    add_rect(s, Inches(0.6), Inches(2.55), Inches(12.1), Inches(2.15),
             fill=SURFACE, line=LINE)
    add_text(s, Inches(0.85), Inches(2.7), Inches(4.0), Inches(0.3),
             "Prompt 1 — 초기 스캐폴딩", size=11, bold=True, color=ACCENT)
    add_text(s, Inches(0.85), Inches(3.0), Inches(11.6), Inches(1.6),
             prompt1, size=10, color=INK, font="Consolas")

    add_rect(s, Inches(0.6), Inches(4.85), Inches(12.1), Inches(2.15),
             fill=SURFACE, line=LINE)
    add_text(s, Inches(0.85), Inches(5.0), Inches(4.0), Inches(0.3),
             "Prompt 2 — FE 3 페이지", size=11, bold=True, color=ACCENT)
    add_text(s, Inches(0.85), Inches(5.3), Inches(11.6), Inches(1.65),
             prompt2, size=10, color=INK, font="Consolas")

    # -------------------------------------------------- 13. ⑥ 바이브코딩 지시 (2/2)
    s = new_slide()
    add_title(s, "바이브코딩 지시 — IaC · 리포트 뷰", eyebrow="Section ⑥  ·  바이브코딩 지시 (2/2)")
    add_section_badge(s, Inches(10.8), Inches(0.6), "⑥ 바이브코딩")

    prompt3 = (
        "역할: DevOps/플랫폼.\n"
        "기능: 프로토타입을 공개 URL 로 배포하기 위한 IaC — Vercel 프로젝트 2개(web/api) + Supabase 프로젝트 1개.\n"
        "기술 스택: Terraform · vercel/vercel ~> 2.0 · supabase/supabase ~> 1.5 · null_resource + local-exec(psql) 로 schema.sql 적용.\n"
        "산출물: infra/terraform/ 에 modules/{vercel,supabase} 구조, `terraform apply` 1회로 전체 배포. "
        "env var: DATABASE_URL · GROQ_API_KEY · NEXT_PUBLIC_API_BASE_URL 자동 주입."
    )
    prompt4 = (
        "역할: 프런트엔드 엔지니어.\n"
        "기능: 주간 리포트 페이지 추가 — API 에서 받은 markdown 을 렌더링 + PDF 로 저장.\n"
        "UI/UX: 서버 컴포넌트로 데이터 페치, 클라이언트 ReportView 가 렌더. "
        "우상단 고정 'Download PDF' 버튼 → window.print() 호출. "
        "@media print 로 네비·버튼 숨김, prose 타이포그래피는 유지.\n"
        "기술 스택: Next.js 16 App Router · react-markdown 10.x · remark-gfm 4.x · tailwindcss/typography(prose)."
    )

    add_rect(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(2.3),
             fill=SURFACE, line=LINE)
    add_text(s, Inches(0.85), Inches(2.25), Inches(5.0), Inches(0.3),
             "Prompt 3 — Terraform IaC", size=11, bold=True, color=ACCENT)
    add_text(s, Inches(0.85), Inches(2.55), Inches(11.6), Inches(1.8),
             prompt3, size=10, color=INK, font="Consolas")

    add_rect(s, Inches(0.6), Inches(4.55), Inches(12.1), Inches(2.3),
             fill=SURFACE, line=LINE)
    add_text(s, Inches(0.85), Inches(4.7), Inches(5.0), Inches(0.3),
             "Prompt 4 — /games/[id]/report", size=11, bold=True, color=ACCENT)
    add_text(s, Inches(0.85), Inches(5.0), Inches(11.6), Inches(1.8),
             prompt4, size=10, color=INK, font="Consolas")

    # -------------------------------------------------- 14. 프로토타입 링크
    s = new_slide()
    add_title(s, "프로토타입 링크", eyebrow="Deliverable ②  ·  Prototype")

    add_rect(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.2), fill=ACCENT_SOFT)
    add_text(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(0.8),
             "Vercel + Supabase 에 Terraform 으로 배포된 실제 동작 프로토타입입니다.\n"
             "아래 URL 에서 내 게임 등록 · 경쟁작 분석 · 주간 리포트 3개 라우트를 체험할 수 있습니다.",
             size=13, color=INK)

    links = [
        ("Web (Frontend)",  PROTOTYPE_WEB_URL, "Next.js 16 · 내 게임 / 경쟁작 / 리포트 3페이지"),
        ("API (Backend)",   PROTOTYPE_API_URL, "FastAPI · read-only + POST /feedback, /games/my"),
        ("Source Code",     GITHUB_REPO_URL,   "전체 구현 · Terraform · 문서 · 스크립트"),
    ]
    top = 3.7
    row_h = 0.95
    for i, (label, url, desc) in enumerate(links):
        y = top + i * row_h
        add_rect(s, Inches(0.6), Inches(y), Inches(2.5), Inches(0.8), fill=ACCENT)
        add_text(s, Inches(0.8), Inches(y + 0.25), Inches(2.3), Inches(0.4),
                 label, size=13, bold=True, color=BG)
        add_rect(s, Inches(3.15), Inches(y), Inches(9.55), Inches(0.8),
                 fill=SURFACE, line=LINE)
        add_text(s, Inches(3.35), Inches(y + 0.1), Inches(9.2), Inches(0.35),
                 url, size=13, bold=True, color=ACCENT, font="Consolas")
        add_text(s, Inches(3.35), Inches(y + 0.45), Inches(9.2), Inches(0.3),
                 desc, size=10, color=MUTE)

    add_text(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.4),
             "※ Supabase 프리 티어는 7일 무활동 시 자동 일시중지 — 심사 기간 동안 주기적으로 접근하여 활성 상태 유지 예정.",
             size=10, color=MUTE)

    # -------------------------------------------------- 15-17. 화면 설계 와이어프레임 3장
    for route_label, route_desc, draw_fn in [
        ("`/` · Home",
         "내 게임 등록 + 게임 카드 그리드. AddMyGameForm 에서 Steam AppID 입력 → POST /games/my → 카드 그리드 업데이트.",
         _draw_wireframe_home),
        ("`/games/[id]` · Competitors",
         "경쟁작 Top 10 + 4축 component score bar + upvote/downvote 이벤트. 우상단 Weekly Report CTA → 리포트 라우트로 이동.",
         _draw_wireframe_competitors),
        ("`/games/[id]/report` · Weekly Report",
         "API 에서 받은 markdown 을 react-markdown+remark-gfm 으로 렌더. 우상단 'Download PDF' → window.print() · @media print 로 UI 숨김.",
         _draw_wireframe_report),
    ]:
        s = new_slide()
        add_title(s, f"화면 설계 · {route_label}", eyebrow="Deliverable ④  ·  Wireframe")
        add_text(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.5),
                 route_desc, size=11, color=MUTE)
        # Wireframe canvas
        draw_fn(s, Inches(0.6), Inches(2.6), Inches(12.1), Inches(4.6))

    # -------------------------------------------------- 18. 데모 영상/GIF
    s = new_slide()
    add_title(s, "데모 영상 · GIF", eyebrow="Deliverable ③  ·  Demo")

    add_rect(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.0), fill=ACCENT_SOFT)
    add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.7),
             "60초 시나리오 · 내 게임 등록부터 주간 리포트 PDF 저장까지 5단계 플로우 전체를 담았습니다.",
             size=13, bold=True, color=INK)

    scenario = [
        ("00:00–00:08", "홈 `/` 로드 → 'My Games' 빈 상태 확인"),
        ("00:08–00:20", "AddMyGameForm 에 Steam appid 입력 (예: 730 = CS2) → 카드 그리드에 추가"),
        ("00:20–00:35", "내 게임 카드 클릭 → `/games/[id]` 경쟁작 Top 10 + component score bar 확인"),
        ("00:35–00:42", "경쟁작 항목에 👍 클릭 → toast 확인 (pm_feedback 기록)"),
        ("00:42–00:55", "상단 'Weekly Report' → `/games/[id]/report` 마크다운 리포트 스크롤"),
        ("00:55–01:00", "우상단 'Download PDF' 클릭 → window.print() 창 확인"),
    ]
    top = 3.4
    row_h = 0.45
    for i, (t, desc) in enumerate(scenario):
        y = top + i * row_h
        if i % 2 == 1:
            add_rect(s, Inches(0.6), Inches(y), Inches(12.1), Inches(row_h), fill=SURFACE)
        add_text(s, Inches(0.85), Inches(y + 0.1), Inches(2.2), Inches(0.3),
                 t, size=12, bold=True, color=ACCENT, font="Consolas")
        add_text(s, Inches(3.2), Inches(y + 0.1), Inches(9.4), Inches(0.3),
                 desc, size=12, color=BODY)

    add_rect(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.75), fill=ACCENT_SOFT)
    add_text(s, Inches(0.85), Inches(6.45), Inches(11.5), Inches(0.5),
             f"GIF 파일: {DEMO_GIF_URL} · 800~1000px 폭 · 15fps · <10MB (ShareX/ScreenToGif)",
             size=12, bold=True, color=INK, font="Consolas")

    # -------------------------------------------------- 19. 부록 - 유사도 공식
    s = new_slide()
    add_title(s, "부록 · 유사도 공식", eyebrow="Appendix")

    add_rect(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.4),
             fill=RGBColor(0x0F, 0x17, 0x2A))
    add_text(s, Inches(0.85), Inches(2.4), Inches(11.5), Inches(1.0),
             "score = 0.40 · semantic + 0.25 · genre + 0.20 · tier + 0.15 · bm",
             size=22, bold=True, color=BG, font="Consolas")
    axes = [
        ("semantic", "cosine(desc_embed_a, desc_embed_b)",
         "description + 리뷰 합본을 sentence-transformers 로 임베딩한 벡터의 코사인 유사도"),
        ("genre",    "cosine(genre_embed_a, genre_embed_b)",
         "LLM 이 추출한 장르 문자열의 임베딩 — 카테고리 이진 매칭이 아닌 연속값"),
        ("tier",     "exp(-|tier_a - tier_b| / 0.3)",
         "log1p(review_count) 정규화 후 거리 기반 Gaussian — 체급 유사도"),
        ("bm",       "1 - KL_sym(bm_a, bm_b) / (2·log 4)",
         "gacha/ads/premium/sub 확률분포의 대칭 KL 거리 → 유사도로 반전"),
    ]
    top = 3.8
    for i, (name, formula, desc) in enumerate(axes):
        y = top + i * 0.7
        add_text(s, Inches(0.6), Inches(y), Inches(1.6), Inches(0.3),
                 name, size=13, bold=True, color=ACCENT, font="Consolas")
        add_text(s, Inches(2.3), Inches(y), Inches(5.0), Inches(0.3),
                 formula, size=11, color=INK, font="Consolas")
        add_text(s, Inches(7.4), Inches(y), Inches(5.3), Inches(0.6),
                 desc, size=11, color=BODY)

    # -------------------------------------------------- 20. 제출 정보 / Thank you
    s = new_slide()
    add_rect(s, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(SLIDE_H_IN), fill=ACCENT)
    add_text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(1.0),
             "Thank you",
             size=60, bold=True, color=BG, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(3.8), Inches(12), Inches(0.5),
             f"{APPLICANT_NAME} · {APPLICANT_POSITION}",
             size=20, color=RGBColor(0xE4, 0xE8, 0xF0), align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(4.4), Inches(12), Inches(0.4),
             f"제출일 · {SUBMISSION_DATE}",
             size=13, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.4),
             f"Source: {GITHUB_REPO_URL}",
             size=11, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER, font="Consolas")

    return slides


# ------------------------------------------------------------------
# Main
# ------------------------------------------------------------------

def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    slides = build(prs)
    total = len(slides)
    # Cover (1) and Thank-you (last) are full-bleed accent → skip footer.
    for i, s in enumerate(slides, start=1):
        if i == 1 or i == total:
            continue
        add_footer(s, i, total)

    out_dir = Path(__file__).resolve().parents[1] / "docs" / "submission"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "submission.pptx"
    prs.save(out_path)
    print(f"[ok] wrote {out_path} ({total} slides)")


if __name__ == "__main__":
    main()
